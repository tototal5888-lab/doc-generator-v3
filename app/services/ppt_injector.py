"""
PPTX Image Injector Service
處理將圖片插入到現有 PowerPoint 簡報的邏輯
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image


class PPTXInjector:
    """PowerPoint 圖片注入器"""
    
    @staticmethod
    def inject_images(pptx_path, injections, output_path=None):
        """
        將圖片注入到指定的 PPTX 文件中
        
        Args:
            pptx_path (str): 源 PPTX 文件路徑
            injections (list): 注入配置列表，格式: [{"image_path": "...", "slide_number": 1}, ...]
            output_path (str, optional): 輸出文件路徑，默認為源文件名 + "_v1"
            
        Returns:
            str: 輸出文件路徑
        """
        # 加載簡報
        prs = Presentation(pptx_path)
        
        # 處理每個注入請求
        for injection in injections:
            image_path = injection.get('image_path')
            slide_number = injection.get('slide_number')
            
            if not image_path or not slide_number:
                continue
                
            # 驗證圖片文件存在
            if not os.path.exists(image_path):
                print(f"警告: 圖片文件不存在: {image_path}")
                continue
            
            # 驗證頁碼有效
            if slide_number < 1 or slide_number > len(prs.slides):
                print(f"警告: 無效的頁碼 {slide_number}，簡報共有 {len(prs.slides)} 頁")
                continue
            
            # 獲取目標 Slide (索引從 0 開始)
            slide = prs.slides[slide_number - 1]
            
            # 插入圖片
            PPTXInjector._insert_image_to_slide(slide, image_path, prs.slide_width, prs.slide_height)
        
        # 生成輸出路徑
        if output_path is None:
            base_name = os.path.splitext(pptx_path)[0]
            output_path = f"{base_name}_v1.pptx"
            
            # 如果 v1 已存在，遞增版本號
            version = 1
            while os.path.exists(output_path):
                version += 1
                output_path = f"{base_name}_v{version}.pptx"
        
        # 保存
        prs.save(output_path)
        return output_path
    
    @staticmethod
    def _insert_image_to_slide(slide, image_path, slide_width, slide_height):
        """
        將圖片插入到 Slide 的右側區域
        
        Args:
            slide: pptx.slide.Slide 對象
            image_path (str): 圖片路徑
            slide_width (int): 簡報寬度
            slide_height (int): 簡報高度
        """
        # 獲取圖片尺寸
        with Image.open(image_path) as img:
            img_width, img_height = img.size
        
        # 計算圖片放置位置和大小
        # 目標：放置在右側 50% 區域，保持比例
        max_width = slide_width * 0.45  # 右側 45% 寬度
        max_height = slide_height * 0.7  # 70% 高度
        
        # 計算縮放比例
        width_ratio = max_width / img_width
        height_ratio = max_height / img_height
        scale_ratio = min(width_ratio, height_ratio)
        
        # 計算實際尺寸
        final_width = int(img_width * scale_ratio)
        final_height = int(img_height * scale_ratio)
        
        # 計算位置（右側居中）
        left = slide_width - final_width - Inches(0.5)  # 右邊距 0.5 英寸
        top = (slide_height - final_height) / 2  # 垂直居中
        
        # 插入圖片
        slide.shapes.add_picture(
            image_path,
            left=left,
            top=top,
            width=final_width,
            height=final_height
        )
