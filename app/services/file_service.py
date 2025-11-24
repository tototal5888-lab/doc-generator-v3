import os
from docx import Document
from pptx import Presentation
try:
    import win32com.client as win32
    import pythoncom
    WIN32_AVAILABLE = True
except ImportError:
    WIN32_AVAILABLE = False

# 可選依賴：PyMuPDF (用於 PDF 讀取)
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

class FileProcessor:
    """文件處理器 - 處理各種格式的文件讀取"""
    
    @staticmethod
    def extract_text_from_pdf(file_path):
        """從PDF提取文本"""
        if not PYMUPDF_AVAILABLE:
            return "錯誤: 未安裝 PyMuPDF，無法讀取 PDF 文件。"
            
        text = ""
        try:
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text()
            return text
        except Exception as e:
            return f"讀取 PDF 失敗: {str(e)}"

    @staticmethod
    def extract_text_from_docx(file_path):
        """從DOCX提取文本"""
        try:
            doc = Document(file_path)
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            return '\n'.join(text)
        except Exception as e:
            return f"讀取 DOCX 失敗: {str(e)}"

    @staticmethod
    def extract_text_from_doc(file_path):
        """從舊版DOC提取文本"""
        if not WIN32_AVAILABLE:
            return "錯誤: 未安裝 pywin32，無法讀取 DOC 文件。"
        try:
            pythoncom.CoInitialize()
            word = win32.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(os.path.abspath(file_path))
            text = doc.Content.Text
            doc.Close()
            word.Quit()
            return text
        except Exception as e:
            return f"讀取 DOC 失敗: {str(e)}"

    @staticmethod
    def extract_text_from_pptx(file_path, include_image_markers=False):
        """
        從PPTX提取文本
        
        Args:
            file_path: PPTX 文件路徑
            include_image_markers: 是否在文本中包含圖片標記
            
        Returns:
            str 或 tuple: 如果 include_image_markers=True，返回 (text, image_count)
        """
        try:
            prs = Presentation(file_path)
            text = []
            total_images = 0
            
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                image_count = 0
                
                for shape in slide.shapes:
                    # 檢查是否為圖片
                    if hasattr(shape, "image"):
                        image_count += 1
                        total_images += 1
                        if include_image_markers:
                            marker = f"\n[圖片 {slide_idx}-{image_count}: 來自投影片 {slide_idx}]\n"
                            slide_text.append(marker)
                    # 提取文字
                    elif hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text.append('\n'.join(slide_text))
            
            result_text = '\n\n'.join(text)
            
            if include_image_markers:
                return result_text, total_images
            else:
                return result_text
                
        except Exception as e:
            if include_image_markers:
                return f"讀取 PPTX 失敗: {str(e)}", 0
            else:
                return f"讀取 PPTX 失敗: {str(e)}"

    @staticmethod
    def extract_text_from_ppt(file_path):
        """從舊版PPT提取文本"""
        if not WIN32_AVAILABLE:
            return "錯誤: 未安裝 pywin32，無法讀取 PPT 文件。"
        try:
            pythoncom.CoInitialize()
            ppt = win32.Dispatch("PowerPoint.Application")
            # ppt.Visible = True # PowerPoint 需要可見才能工作
            presentation = ppt.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
            text = []
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        if shape.TextFrame.HasText:
                            text.append(shape.TextFrame.TextRange.Text)
            presentation.Close()
            ppt.Quit()
            return '\n'.join(text)
        except Exception as e:
            return f"讀取 PPT 失敗: {str(e)}"

    @staticmethod
    def extract_text(file_path):
        """根據文件類型提取文本"""
        ext = file_path.lower().split('.')[-1]
        
        if ext == 'pdf':
            return FileProcessor.extract_text_from_pdf(file_path)
        elif ext == 'docx':
            return FileProcessor.extract_text_from_docx(file_path)
        elif ext == 'doc':
            return FileProcessor.extract_text_from_doc(file_path)
        elif ext == 'pptx':
            return FileProcessor.extract_text_from_pptx(file_path)
        elif ext == 'ppt':
            return FileProcessor.extract_text_from_ppt(file_path)
        elif ext == 'txt' or ext == 'md':
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                try:
                    with open(file_path, 'r', encoding='big5') as f:
                        return f.read()
                except:
                    return "錯誤: 無法識別文件編碼"
        else:
            return "不支持的文件格式"
