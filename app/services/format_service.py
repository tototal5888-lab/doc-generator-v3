import os
import re
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

class FormatConverter:
    """格式轉換器 - 將 Markdown 轉換為各種輸出格式"""

    @staticmethod
    def markdown_to_docx(content, doc_config, template_path=None):
        """將Markdown轉換為DOCX
        
        Args:
            content: Markdown 內容
            doc_config: 文檔配置
            template_path: 模板文件路徑（可選），如果是 DOCX 會繼承其樣式
        """
        # 如果模板是 DOCX，使用它作為基底（繼承樣式和背景）
        if template_path and template_path.lower().endswith('.docx'):
            try:
                doc = Document(template_path)
                # 移除模板原有內容，只保留樣式
                for element in list(doc.element.body):
                    doc.element.body.remove(element)
                print(f"[INFO] 使用 DOCX 模板作為基底: {template_path}")
            except Exception as e:
                print(f"[WARNING] 無法使用 DOCX 模板，使用空白文檔: {e}")
                doc = Document()
        else:
            doc = Document()
        
        # 添加標題
        doc.add_heading(doc_config.get('title', '生成的文檔'), 0)
        
        # 處理內容
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('# '):
                doc.add_heading(line[2:], 1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], 2)
            elif line.startswith('### '):
                doc.add_heading(line[4:], 3)
            elif line.startswith('- ') or line.startswith('* '):
                doc.add_paragraph(line[2:], style='List Bullet')
            elif re.match(r'^\d+\. ', line):
                doc.add_paragraph(line.split('. ', 1)[1], style='List Number')
            else:
                doc.add_paragraph(line)
                
        return doc

    @staticmethod
    def markdown_to_pptx(content, doc_config, image_folder=None, template_path=None):
        """
        將Markdown轉換為PPTX
        
        Args:
            content: Markdown 內容
            doc_config: 文檔配置
            image_folder: 圖片文件夾路徑（可選）
            template_path: 模板文件路徑（可選），如果是 PPTX 會繼承其母片樣式
        """
        # 如果模板是 PPTX，使用它作為基底（繼承母片樣式和背景）
        if template_path and template_path.lower().endswith('.pptx'):
            try:
                prs = Presentation(template_path)
                # 刪除模板原有投影片，只保留母片樣式
                while len(prs.slides) > 0:
                    rId = prs.slides._sldIdLst[0].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[0]
                print(f"[INFO] 使用 PPTX 模板作為基底: {template_path}")
            except Exception as e:
                print(f"[WARNING] 無法使用 PPTX 模板，使用空白簡報: {e}")
                prs = Presentation()
        else:
            prs = Presentation()
        
        # 智能選擇版面配置
        # 需要根據名稱判斷哪個是真正的內容版面
        title_slide_layout = None
        bullet_slide_layout = None
        blank_slide_layout = None
        
        # 先列出所有版面以便除錯
        print(f"[DEBUG] 模板共有 {len(prs.slide_layouts)} 個版面:")
        for i, layout in enumerate(prs.slide_layouts):
            print(f"  [{i}] {layout.name}")
        
        # 遍歷版面，智能選擇
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = layout.name.lower() if layout.name else ""
            layout_name_orig = layout.name if layout.name else ""
            
            # 標題版面（封面）：通常索引0，且名稱只有「標題」沒有「物件」或「內容」
            if title_slide_layout is None:
                if i == 0:  # 第一個通常是封面
                    title_slide_layout = layout
                elif "title" in layout_name and "content" not in layout_name and "物件" not in layout_name_orig:
                    title_slide_layout = layout
            
            # 內容版面：名稱包含「物件」、「內容」、「content」、「body」
            if bullet_slide_layout is None:
                if "物件" in layout_name_orig or "內容" in layout_name_orig:
                    bullet_slide_layout = layout
                    print(f"[DEBUG] 找到內容版面 (中文): [{i}] {layout.name}")
                elif "content" in layout_name or "body" in layout_name:
                    bullet_slide_layout = layout
                    print(f"[DEBUG] 找到內容版面 (英文): [{i}] {layout.name}")
            
            # 空白版面
            if blank_slide_layout is None:
                if "blank" in layout_name or "空白" in layout_name_orig:
                    blank_slide_layout = layout
        
        # 如果找不到「內容版面」，嘗試使用索引 1 或 2
        if bullet_slide_layout is None:
            if len(prs.slide_layouts) > 2:
                bullet_slide_layout = prs.slide_layouts[2]  # 很多模板第三個是內容版面
            elif len(prs.slide_layouts) > 1:
                bullet_slide_layout = prs.slide_layouts[1]
            elif len(prs.slide_layouts) > 0:
                bullet_slide_layout = prs.slide_layouts[0]
        
        # 如果找不到標題版面
        if title_slide_layout is None and len(prs.slide_layouts) > 0:
            title_slide_layout = prs.slide_layouts[0]
        
        if blank_slide_layout is None:
            blank_slide_layout = prs.slide_layouts[0]
        
        # 如果完全沒有版面，使用空白簡報
        if title_slide_layout is None:
            print("[WARNING] 模板沒有任何版面配置，使用空白簡報")
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            bullet_slide_layout = prs.slide_layouts[1]
            blank_slide_layout = prs.slide_layouts[6]
        
        # 重要：如果「標題及物件」版面沒有內容區域，使用空白版面代替
        # 這樣我們可以完全控制 TextBox 的位置
        if bullet_slide_layout and blank_slide_layout:
            print(f"[INFO] 內容頁將使用空白版面 '{blank_slide_layout.name}'，並自動創建 TextBox")
            bullet_slide_layout = blank_slide_layout  # 使用空白版面作為內容版面
        
        print(f"[INFO] 標題版面: {title_slide_layout.name}, 內容版面: {bullet_slide_layout.name if bullet_slide_layout else 'None'}")
        
        # 創建標題頁
        slide = prs.slides.add_slide(title_slide_layout)
        
        # 設定標題（如果有的話）
        if slide.shapes.title:
            slide.shapes.title.text = doc_config.get('title', '生成的演示文稿')
        
        # 設定副標題（如果有的話）
        try:
            if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
                slide.placeholders[1].text = "由 AI 自動生成"
        except Exception:
            pass  # 沒有副標題 placeholder 就跳過
        
        current_slide = None
        body_shape = None
        tf = None
        
        # 圖片標記的正則表達式：支持多種格式
        # 格式1: [圖片 1-1: 來自投影片 1]
        # 格式2: - 圖片 1-1: 來自投影片 1
        image_pattern = re.compile(r'[-\[]?\s*圖片\s+(\d+)-(\d+)(?::\s*來自投影片\s*\d+)?[\]]?')
        
        lines = content.split('\n')
        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                continue
            
            # 檢查是否為圖片標記
            image_match = image_pattern.search(line_stripped)
            if image_match and image_folder:
                slide_num = image_match.group(1)
                img_num = image_match.group(2)
                
                # 構建圖片文件名
                image_filename = f"slide_{slide_num}_image_{img_num}"
                
                # 查找匹配的圖片文件（可能有不同擴展名）
                image_path = None
                for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
                    potential_path = os.path.join(image_folder, image_filename + ext)
                    if os.path.exists(potential_path):
                        image_path = potential_path
                        break
                
                if image_path:
                    # 如果當前沒有投影片，創建一個帶標題的投影片
                    if not current_slide:
                        current_slide = prs.slides.add_slide(bullet_slide_layout)
                        shapes = current_slide.shapes
                        if shapes.title:
                            shapes.title.text = "圖片內容"
                        # 安全取得 body_shape
                        body_shape = None
                        try:
                            if 1 in [p.placeholder_format.idx for p in current_slide.placeholders]:
                                body_shape = current_slide.placeholders[1]
                                tf = body_shape.text_frame
                        except Exception:
                            tf = None
                    
                    # 在當前投影片中插入圖片（放在右側或下方）
                    left = Inches(5.5)  # 靠右放置
                    top = Inches(1.5)
                    width = Inches(3.5)  # 設定寬度，高度自動調整
                    
                    try:
                        current_slide.shapes.add_picture(image_path, left, top, width=width)
                        print(f"成功插入圖片: {image_path}")
                    except Exception as e:
                        print(f"插入圖片失敗: {image_path}, 錯誤: {str(e)}")
                else:
                    print(f"未找到圖片文件: {image_filename} (在 {image_folder})")
                
                continue
                
            # 新的一頁（一級或二級標題）
            if line_stripped.startswith('# ') or line_stripped.startswith('## '):
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                
                # 取得標題文字
                title_text = line.lstrip('#').strip()
                
                # 檢查版面是否有標題 placeholder
                if shapes.title:
                    shapes.title.text = title_text
                    print(f"[DEBUG] 使用 placeholder 設定標題")
                else:
                    # 空白版面沒有標題 placeholder，創建標題 TextBox
                    print(f"[DEBUG] 空白版面，創建標題 TextBox")
                    title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                    title_box.text_frame.paragraphs[0].text = title_text
                    # 設定標題樣式（較大字體）
                    from pptx.util import Pt
                    title_box.text_frame.paragraphs[0].font.size = Pt(28)
                    title_box.text_frame.paragraphs[0].font.bold = True
                
                # 除錯：列出所有 placeholders
                print(f"[DEBUG] Slide placeholders: {len(list(current_slide.placeholders))} 個")
                
                # 對於空白版面，直接創建內容 TextBox
                print(f"[INFO] 創建內容 TextBox")
                # 在標題下方創建一個 TextBox
                left = Inches(0.5)
                top = Inches(1.5)  # 標題下方
                width = Inches(9)
                height = Inches(5.5)
                textbox = shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                
            # 內容點
            elif current_slide and (line.startswith('- ') or line.startswith('* ')):
                if tf:
                    p = tf.add_paragraph()
                    p.text = line[2:]
                    p.level = 0
            
            # 普通文本作為內容
            elif current_slide and not line.startswith('#'):
                if tf:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
                    
        return prs

    @staticmethod
    def markdown_to_pdf(content, doc_config, output_path):
        """將Markdown轉換為PDF"""
        
        # 註冊中文字體 (如果有的話，否則使用默認)
        # 這裡假設系統有微軟正黑體，如果沒有可能需要調整
        try:
            pdfmetrics.registerFont(TTFont('MsJhengHei', 'msjh.ttc'))
            font_name = 'MsJhengHei'
        except:
            try:
                pdfmetrics.registerFont(TTFont('ArialUnicode', 'arialuni.ttf'))
                font_name = 'ArialUnicode'
            except:
                font_name = 'Helvetica' #  fallback
        
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # 定義樣式
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontName=font_name,
            fontSize=24,
            spaceAfter=30
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading1'],
            fontName=font_name,
            fontSize=18,
            spaceBefore=20,
            spaceAfter=10
        )
        
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontName=font_name,
            fontSize=12,
            spaceBefore=6,
            leading=18
        )
        
        story = []
        
        # 添加標題
        story.append(Paragraph(doc_config.get('title', '生成的文檔'), title_style))
        story.append(Spacer(1, 12))
        
        # 處理內容
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                story.append(Spacer(1, 6))
                continue
                
            if line.startswith('# ') or line.startswith('## '):
                text = line.lstrip('#').strip()
                story.append(Paragraph(text, heading_style))
            elif line.startswith('- ') or line.startswith('* '):
                text = "• " + line[2:]
                story.append(Paragraph(text, normal_style))
            else:
                story.append(Paragraph(line, normal_style))
                
        doc.build(story)
        return output_path
