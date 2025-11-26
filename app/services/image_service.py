import os
import uuid
from pptx import Presentation
from PIL import Image
import io

class ImageExtractor:
    """圖片提取器 - 從 PPTX 文件中提取圖片"""
    
    @staticmethod
    def extract_images_from_pptx(pptx_path, output_folder):
        """
        從 PPTX 提取所有圖片
        
        Args:
            pptx_path: PPTX 文件路徑
            output_folder: 圖片輸出文件夾
            
        Returns:
            list: 圖片信息列表 [{'slide': 1, 'index': 1, 'path': 'xxx.png', 'marker': '[圖片 1-1]'}]
        """
        try:
            prs = Presentation(pptx_path)
            images = []
            
            # 確保輸出文件夾存在
            os.makedirs(output_folder, exist_ok=True)
            
            # 遍歷每個投影片
            for slide_idx, slide in enumerate(prs.slides, start=1):
                image_count = 0
                
                # 遍歷投影片中的所有形狀
                for shape in slide.shapes:
                    # 檢查是否為圖片
                    if hasattr(shape, "image"):
                        image_count += 1
                        
                        # 獲取圖片數據
                        image = shape.image
                        image_bytes = image.blob
                        
                        # 確定圖片格式
                        ext = image.ext
                        if not ext.startswith('.'):
                            ext = '.' + ext
                        
                        # 生成圖片文件名
                        filename = f"slide_{slide_idx}_image_{image_count}{ext}"
                        filepath = os.path.join(output_folder, filename)
                        
                        # 保存圖片
                        with open(filepath, 'wb') as f:
                            f.write(image_bytes)
                        
                        # 記錄圖片信息
                        marker = f"[圖片 {slide_idx}-{image_count}]"
                        images.append({
                            'slide': slide_idx,
                            'index': image_count,
                            'path': filepath,
                            'filename': filename,
                            'marker': marker,
                            'ext': ext
                        })
            
            return images
            
        except Exception as e:
            raise Exception(f"提取圖片失敗: {str(e)}")
    
    @staticmethod
    def create_temp_image_folder(base_folder='temp_images'):
        """
        創建臨時圖片文件夾
        
        Args:
            base_folder: 基礎文件夾名稱
            
        Returns:
            str: 臨時文件夾路徑
        """
        # 使用 UUID 創建唯一的臨時文件夾
        session_id = str(uuid.uuid4())[:8]
        temp_folder = os.path.join(base_folder, session_id)
        os.makedirs(temp_folder, exist_ok=True)
        return temp_folder
    
    @staticmethod
    def cleanup_temp_images(temp_folder):
        """
        清理臨時圖片文件夾
        
        Args:
            temp_folder: 臨時文件夾路徑
        """
        try:
            if os.path.exists(temp_folder):
                import shutil
                shutil.rmtree(temp_folder)
        except Exception as e:
            print(f"清理臨時文件夾失敗: {str(e)}")
