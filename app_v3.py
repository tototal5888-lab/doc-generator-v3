"""
AI文檔生成器 V3.0 - 多格式輸出版本（完整終極版）
支持輸出格式: DOCX, PPTX, PDF, Markdown
支持 AI: Google Gemini / OpenAI GPT（帶成本追蹤）
可選依賴：PyMuPDF (PDF 讀取), pywin32 (舊版 Office)
Author: Jim
Date: 2025-11-19
Version: 3.0.3 (終極版 - 含成本追蹤)
"""

import os
import json
import logging
from datetime import datetime
from pathlib import Path
from flask import Flask, request, jsonify, render_template, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
import requests
from typing import Dict, List, Optional, Union
import tempfile
import base64
import subprocess
import markdown
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import time

# 可選依賴：PyMuPDF (用於 PDF 讀取)
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("⚠️  警告：PyMuPDF 未安裝，PDF 模板讀取功能將不可用")
    print("    建議：使用 Word (.docx) 或文本 (.txt) 格式的模板")

# 可選依賴：pywin32 (用於讀取舊版 Office 文件)
try:
    import win32com.client
    import pythoncom
    PYWIN32_AVAILABLE = True
except ImportError:
    PYWIN32_AVAILABLE = False
    print("⚠️  警告：pywin32 未安裝，舊版 .doc/.ppt 讀取功能將不可用")

# 配置日誌
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Flask 應用配置
app = Flask(__name__)
app.config['SECRET_KEY'] = 'your-secret-key-change-in-production'
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['TEMPLATES_FOLDER'] = 'templates_storage'
app.config['OUTPUT_FOLDER'] = 'output'

# 支持的文件類型
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'doc', 'pptx', 'ppt', 'txt', 'md'}
OUTPUT_FORMATS = ['docx', 'pptx', 'pdf', 'md']

# OpenAI 價格配置（每 1M tokens 的價格，單位：美元）
OPENAI_PRICING = {
    'gpt-4o': {
        'input': 2.50,
        'output': 10.00
    },
    'gpt-4o-mini': {
        'input': 0.15,
        'output': 0.60
    },
    'gpt-4-turbo': {
        'input': 10.00,
        'output': 30.00
    },
    'gpt-4': {
        'input': 30.00,
        'output': 60.00
    },
    'gpt-3.5-turbo': {
        'input': 0.50,
        'output': 1.50
    }
}

# 啟用 CORS
CORS(app, resources={r"/api/*": {"origins": "*"}})

# 創建必要的目錄
for folder in ['uploads', 'templates_storage', 'output', 'config']:
    os.makedirs(folder, exist_ok=True)

# 文檔類型配置
DOCUMENT_TYPES = {
    "system_doc": {
        "name": "系統文檔",
        "description": "系統設計、架構、技術規格等文檔",
        "prompt_template": """
請根據以下模板和要求生成系統文檔：

文檔類型：系統文檔
模板內容：
{template_content}

要求：
1. 保持專業的技術文檔風格
2. 包含系統架構、功能模組、技術棧等內容
3. 確保文檔結構清晰、邏輯嚴謹
4. 使用標準的技術術語
5. 根據模板格式調整輸出格式

用戶需求：
{user_requirements}

請生成完整的系統文檔內容，使用Markdown格式輸出。
        """,
        "sections": ["系統概述", "架構設計", "功能模組", "技術棧", "部署說明", "維護指南"]
    },
    "sop": {
        "name": "SOP標準作業程序",
        "description": "標準作業程序文檔",
        "prompt_template": """
請根據以下模板和要求生成SOP文檔：

文檔類型：SOP標準作業程序
模板內容：
{template_content}

要求：
1. 步驟清晰、具體、可執行
2. 包含目的、範圍、職責、程序步驟等
3. 注意事項和異常處理要完整
4. 符合ISO標準格式
5. 根據模板格式調整輸出格式

用戶需求：
{user_requirements}

請生成完整的SOP文檔內容，使用Markdown格式輸出。
        """,
        "sections": ["目的", "適用範圍", "職責分工", "作業程序", "注意事項", "異常處理", "相關文件"]
    },
    "technical_report": {
        "name": "技術報告",
        "description": "技術分析、研究報告、評估文檔",
        "prompt_template": """
請根據以下模板和要求生成技術報告：

文檔類型：技術報告
模板內容：
{template_content}

要求：
1. 數據準確、分析深入
2. 包含背景、方法、結果、結論等部分
3. 圖表說明清晰
4. 技術細節完整
5. 根據模板格式調整輸出格式

用戶需求：
{user_requirements}

請生成完整的技術報告內容，使用Markdown格式輸出。
        """,
        "sections": ["執行摘要", "背景介紹", "技術分析", "實施方案", "結果評估", "結論建議"]
    }
}


class FileProcessor:
    """文件處理器 - 處理各種格式的文件讀取"""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """從PDF提取文本"""
        if not PYMUPDF_AVAILABLE:
            logger.warning("PyMuPDF 未安裝，無法提取 PDF 文本")
            return """# PDF 模板

⚠️ 注意：PDF 讀取功能需要安裝 PyMuPDF 套件。

建議：
- 使用 Word (.docx) 格式的模板
- 或使用純文本 (.txt) 格式的模板
- 或執行：pip install PyMuPDF

如需使用 PDF 模板，請安裝後重啟應用。
"""
        
        try:
            text = []
            pdf_document = fitz.open(file_path)
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                text.append(page.get_text())
            pdf_document.close()
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"PDF提取失敗: {e}")
            return f"# PDF 讀取錯誤\n\n無法讀取 PDF 文件：{str(e)}"
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """從DOCX提取文本"""
        try:
            doc = Document(file_path)
            text = []
            
            # 提取段落
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            # 提取表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text.append(' | '.join(row_text))
            
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"DOCX提取失敗: {e}")
            return f"# DOCX 讀取錯誤\n\n無法讀取 Word 文件：{str(e)}"
    
    @staticmethod
    def extract_text_from_doc(file_path: str) -> str:
        """從舊版DOC提取文本"""
        if not PYWIN32_AVAILABLE:
            logger.warning("pywin32 未安裝，嘗試作為 DOCX 處理")
            try:
                return FileProcessor.extract_text_from_docx(file_path)
            except:
                return "# 舊版 DOC 文件\n\n無法讀取舊版 .doc 文件，請轉換為 .docx 格式。"
        
        try:
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(os.path.abspath(file_path))
            text = doc.Range().Text
            doc.Close()
            word.Quit()
            pythoncom.CoUninitialize()
            return text
        except Exception as e:
            logger.error(f"DOC提取失敗: {e}")
            try:
                return FileProcessor.extract_text_from_docx(file_path)
            except:
                return f"# DOC 讀取錯誤\n\n無法讀取文件：{str(e)}"
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """從PPTX提取文本"""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"\n--- Slide {slide_num} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        if shape.text.strip():
                            text.append(shape.text)
            
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"PPTX提取失敗: {e}")
            return f"# PPTX 讀取錯誤\n\n無法讀取 PowerPoint 文件：{str(e)}"
    
    @staticmethod
    def extract_text_from_ppt(file_path: str) -> str:
        """從舊版PPT提取文本"""
        if not PYWIN32_AVAILABLE:
            logger.warning("pywin32 未安裝，嘗試作為 PPTX 處理")
            try:
                return FileProcessor.extract_text_from_pptx(file_path)
            except:
                return "# 舊版 PPT 文件\n\n無法讀取舊版 .ppt 文件，請轉換為 .pptx 格式。"
        
        try:
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False
            presentation = powerpoint.Presentations.Open(os.path.abspath(file_path))
            text = []
            
            for slide in presentation.Slides:
                text.append(f"\n--- Slide {slide.SlideIndex} ---\n")
                for shape in slide.Shapes:
                    if hasattr(shape, 'TextFrame'):
                        if shape.TextFrame.HasText:
                            text.append(shape.TextFrame.TextRange.Text)
            
            presentation.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"PPT提取失敗: {e}")
            try:
                return FileProcessor.extract_text_from_pptx(file_path)
            except:
                return f"# PPT 讀取錯誤\n\n無法讀取文件：{str(e)}"
    
    @staticmethod
    def extract_text(file_path: str) -> str:
        """根據文件類型提取文本"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return FileProcessor.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return FileProcessor.extract_text_from_docx(file_path)
        elif ext == '.doc':
            return FileProcessor.extract_text_from_doc(file_path)
        elif ext == '.pptx':
            return FileProcessor.extract_text_from_pptx(file_path)
        elif ext == '.ppt':
            return FileProcessor.extract_text_from_ppt(file_path)
        elif ext in ['.txt', '.md']:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"文本文件讀取失敗: {e}")
                return f"# 文本讀取錯誤\n\n無法讀取文件：{str(e)}"
        
        return "# 不支持的文件格式\n\n請使用 DOCX, PPTX, TXT 或 MD 格式的文件。"


class FormatConverter:
    """格式轉換器 - 將 Markdown 轉換為各種輸出格式"""
    
    @staticmethod
    def markdown_to_docx(content: str, doc_config: dict) -> Document:
        """將Markdown轉換為DOCX"""
        doc = Document()
        
        # 添加標題
        title = doc.add_heading(doc_config['name'], level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 添加元信息
        meta = doc.add_paragraph()
        meta.add_run(f"生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n").font.size = Pt(10)
        meta.add_run(f"文檔類型：{doc_config['name']}").font.size = Pt(10)
        doc.add_paragraph("")
        
        # 解析 Markdown 內容
        lines = content.split('\n')
        in_code_block = False
        
        for line in lines:
            # 處理代碼塊
            if line.strip().startswith('```'):
                in_code_block = not in_code_block
                continue
            
            if in_code_block:
                p = doc.add_paragraph(line, style='Normal')
                p.paragraph_format.left_indent = Inches(0.5)
                run = p.runs[0]
                run.font.name = 'Courier New'
                run.font.size = Pt(9)
                continue
            
            # 處理正常內容
            if line.strip():
                if line.startswith('###'):
                    doc.add_heading(line.replace('#', '').strip(), level=3)
                elif line.startswith('##'):
                    doc.add_heading(line.replace('#', '').strip(), level=2)
                elif line.startswith('#'):
                    doc.add_heading(line.replace('#', '').strip(), level=1)
                elif line.strip().startswith('- ') or line.strip().startswith('* '):
                    doc.add_paragraph(line.strip()[2:], style='List Bullet')
                elif line.strip()[0:1].isdigit() and '. ' in line:
                    doc.add_paragraph(line.strip().split('. ', 1)[1], style='List Number')
                else:
                    doc.add_paragraph(line.strip())
        
        return doc
    
    @staticmethod
    def markdown_to_pptx(content: str, doc_config: dict) -> Presentation:
        """將Markdown轉換為PPTX"""
        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # 標題頁
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = doc_config['name']
        subtitle.text = f"生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # 解析內容為章節
        sections = []
        current_section = {'title': '', 'content': []}
        
        lines = content.split('\n')
        for line in lines:
            if line.startswith('##'):
                if current_section['title']:
                    sections.append(current_section)
                current_section = {'title': line.replace('#', '').strip(), 'content': []}
            elif line.strip() and not line.startswith('#'):
                current_section['content'].append(line.strip())
        
        if current_section['title']:
            sections.append(current_section)
        
        # 為每個章節創建投影片
        for section in sections:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = section['title']
            
            tf = body_shape.text_frame
            tf.clear()
            
            # 添加內容（最多7條）
            for item in section['content'][:7]:
                if item.startswith('- ') or item.startswith('* '):
                    p = tf.add_paragraph()
                    p.text = item[2:]
                    p.level = 0
                else:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
        
        return prs
    
    @staticmethod
    def markdown_to_pdf(content: str, doc_config: dict, output_path: str):
        """將Markdown轉換為PDF"""
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        
        # 自定義樣式
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor='#2c3e50',
            spaceAfter=30,
            alignment=1
        )
        
        heading1_style = ParagraphStyle(
            'CustomHeading1',
            parent=styles['Heading1'],
            fontSize=16,
            textColor='#34495e',
            spaceAfter=12
        )
        
        story = []
        
        # 添加標題
        story.append(Paragraph(doc_config['name'], title_style))
        story.append(Spacer(1, 12))
        
        # 添加元信息
        meta_text = f"生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        story.append(Paragraph(meta_text, styles['Normal']))
        story.append(Spacer(1, 20))
        
        # 解析內容
        lines = content.split('\n')
        for line in lines:
            if line.strip():
                if line.startswith('##'):
                    story.append(Spacer(1, 12))
                    story.append(Paragraph(line.replace('#', '').strip(), heading1_style))
                elif line.startswith('#'):
                    story.append(PageBreak())
                    story.append(Paragraph(line.replace('#', '').strip(), title_style))
                else:
                    text = line.strip()
                    if text.startswith('- ') or text.startswith('* '):
                        text = '• ' + text[2:]
                    story.append(Paragraph(text, styles['Normal']))
                    story.append(Spacer(1, 6))
        
        doc.build(story)


class DocumentGenerator:
    """文檔生成器 - 核心業務邏輯"""
    
    def __init__(self):
        self.api_config = self.load_api_config()
        self.api_type = self.api_config.get('api_type', 'gemini')
    
    def load_api_config(self) -> dict:
        """加載 API 配置"""
        config_path = 'config/api_config.json'
        if os.path.exists(config_path):
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception as e:
                logger.error(f"加載配置失敗: {e}")
        return {}
    
    def save_api_config(self, config: dict):
        """保存 API 配置"""
        config_path = 'config/api_config.json'
        try:
            with open(config_path, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
            self.api_config = config
            self.api_type = config.get('api_type', 'gemini')
            logger.info("API 配置已保存")
        except Exception as e:
            logger.error(f"保存配置失敗: {e}")
            raise
    
    def log_cost_to_file(self, model: str, input_tokens: int, output_tokens: int, cost: float):
        """記錄成本到文件"""
        cost_log_file = 'config/cost_log.csv'
        
        # 如果文件不存在，創建並寫入表頭
        if not os.path.exists(cost_log_file):
            with open(cost_log_file, 'w', encoding='utf-8') as f:
                f.write('時間,模型,輸入Tokens,輸出Tokens,總Tokens,成本(美元),成本(台幣)\n')
        
        # 追加記錄
        with open(cost_log_file, 'a', encoding='utf-8') as f:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            total_tokens = input_tokens + output_tokens
            cost_twd = cost * 32
            f.write(f'{timestamp},{model},{input_tokens},{output_tokens},{total_tokens},{cost:.6f},{cost_twd:.4f}\n')
    
    def call_gemini_api(self, prompt: str) -> str:
        """調用 Gemini API（帶重試機制）"""
        api_key = self.api_config.get('gemini_api_key')
        if not api_key:
            raise ValueError("未設置 Gemini API 密鑰，請在設定頁面配置")
        
        # Gemini API 端點
        url = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash-exp:generateContent"
        
        headers = {'Content-Type': 'application/json'}
        params = {'key': api_key}
        
        data = {
            "contents": [{
                "parts": [{"text": prompt}]
            }],
            "generationConfig": {
                "temperature": 0.7,
                "topP": 0.8,
                "topK": 40,
                "maxOutputTokens": 8192
            }
        }
        
        # 重試配置
        max_retries = 3
        base_delay = 15  # 基礎延遲（秒）
        
        for attempt in range(max_retries):
            try:
                logger.info(f"調用 Gemini API... (嘗試 {attempt + 1}/{max_retries})")
                
                response = requests.post(
                    url, 
                    headers=headers, 
                    params=params, 
                    json=data, 
                    timeout=120
                )
                
                # 檢查響應狀態
                if response.status_code == 200:
                    result = response.json()
                    
                    # 檢查是否有有效內容
                    if 'candidates' in result and len(result['candidates']) > 0:
                        content = result['candidates'][0]['content']['parts'][0]['text']
                        logger.info("✅ API 調用成功")
                        return content
                    else:
                        raise ValueError("API 返回了空內容")
                
                elif response.status_code == 429:  # 請求頻率限制
                    if attempt < max_retries - 1:
                        # 指數退避策略
                        wait_time = base_delay * (2 ** attempt)  # 15s, 30s, 60s
                        logger.warning(
                            f"⚠️  API 請求頻率限制 (429)，等待 {wait_time} 秒後重試..."
                            f" (第 {attempt + 1}/{max_retries} 次嘗試)"
                        )
                        time.sleep(wait_time)
                        continue
                    else:
                        raise ValueError(
                            "API 請求頻率限制，已達到最大重試次數。\n\n"
                            "建議：\n"
                            "1. 等待 3-5 分鐘後再試\n"
                            "2. 檢查今日配額：https://aistudio.google.com/app/apikey\n"
                            "3. 控制請求頻率，每次生成間隔至少 15 秒\n"
                            "4. 或切換到 OpenAI API"
                        )
                
                elif response.status_code == 400:  # 錯誤請求
                    error_msg = response.json().get('error', {}).get('message', '未知錯誤')
                    raise ValueError(f"API 請求錯誤 (400): {error_msg}")
                
                elif response.status_code == 403:  # 權限問題
                    raise ValueError(
                        "API 權限錯誤 (403)。\n\n"
                        "可能原因：\n"
                        "1. API Key 無效或已過期\n"
                        "2. API 未啟用\n"
                        "請檢查：https://aistudio.google.com/app/apikey"
                    )
                
                elif response.status_code == 404:  # 找不到資源
                    raise ValueError(
                        "API 端點錯誤 (404)。\n"
                        "模型可能不可用或 URL 不正確。"
                    )
                
                elif response.status_code >= 500:  # 服務器錯誤
                    if attempt < max_retries - 1:
                        wait_time = base_delay * (attempt + 1)
                        logger.warning(f"服務器錯誤 ({response.status_code})，等待 {wait_time} 秒後重試...")
                        time.sleep(wait_time)
                        continue
                    else:
                        raise ValueError(f"Gemini 服務器錯誤 ({response.status_code})，請稍後再試")
                
                else:
                    # 其他錯誤
                    response.raise_for_status()
            
            except requests.exceptions.Timeout:
                if attempt < max_retries - 1:
                    logger.warning(f"請求超時，重試中... ({attempt + 1}/{max_retries})")
                    time.sleep(10)
                    continue
                else:
                    raise ValueError("API 請求超時，請檢查網絡連接")
            
            except requests.exceptions.RequestException as e:
                raise ValueError(f"網絡請求失敗：{str(e)}")
            
            except KeyError as e:
                raise ValueError(f"API 響應格式錯誤：缺少字段 {str(e)}")
            
            except Exception as e:
                if "429" in str(e) and attempt < max_retries - 1:
                    wait_time = base_delay * (2 ** attempt)
                    logger.warning(f"遇到錯誤，等待 {wait_time} 秒後重試...")
                    time.sleep(wait_time)
                    continue
                raise ValueError(f"API 調用失敗：{str(e)}")
        
        raise ValueError("API 調用失敗：已達到最大重試次數")
    
    def call_openai_api(self, prompt: str) -> str:
        """調用 OpenAI API（帶成本追蹤）"""
        api_key = self.api_config.get('openai_api_key')
        model = self.api_config.get('openai_model', 'gpt-4o-mini')
        
        if not api_key:
            raise ValueError("未設置 OpenAI API 密鑰，請在設定頁面配置")
        
        url = "https://api.openai.com/v1/chat/completions"
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {api_key}'
        }
        data = {
            "model": model,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.7,
            "max_tokens": 8000
        }
        
        try:
            logger.info(f"調用 OpenAI API ({model})...")
            response = requests.post(url, headers=headers, json=data, timeout=120)
            response.raise_for_status()
            result = response.json()
            
            # 提取內容
            content = result['choices'][0]['message']['content']
            
            # 提取使用量
            usage = result.get('usage', {})
            input_tokens = usage.get('prompt_tokens', 0)
            output_tokens = usage.get('completion_tokens', 0)
            total_tokens = usage.get('total_tokens', 0)
            
            # 計算成本
            pricing = OPENAI_PRICING.get(model, OPENAI_PRICING['gpt-4o-mini'])
            input_cost = (input_tokens / 1_000_000) * pricing['input']
            output_cost = (output_tokens / 1_000_000) * pricing['output']
            total_cost = input_cost + output_cost
            
            # 記錄使用情況
            logger.info("✅ API 調用成功")
            logger.info(f"📊 Token 使用量:")
            logger.info(f"   - 輸入: {input_tokens:,} tokens")
            logger.info(f"   - 輸出: {output_tokens:,} tokens")
            logger.info(f"   - 總計: {total_tokens:,} tokens")
            logger.info(f"💰 本次成本:")
            logger.info(f"   - 輸入成本: ${input_cost:.6f}")
            logger.info(f"   - 輸出成本: ${output_cost:.6f}")
            logger.info(f"   - 總成本: ${total_cost:.6f} (約 NT$ {total_cost * 32:.2f})")
            
            # 記錄到文件
            self.log_cost_to_file(model, input_tokens, output_tokens, total_cost)
            
            return content
            
        except requests.exceptions.HTTPError as e:
            error_detail = ""
            try:
                error_detail = e.response.json()
                error_msg = error_detail.get('error', {}).get('message', '未知錯誤')
                logger.error(f"詳細錯誤：{error_msg}")
                
                # 提供更友好的錯誤提示
                if '401' in str(e) or 'invalid' in error_msg.lower():
                    raise ValueError(
                        "OpenAI API Key 無效。\n\n"
                        "請檢查：\n"
                        "1. API Key 格式是否正確 (sk-proj-... 或 sk-...)\n"
                        "2. 訪問 https://platform.openai.com/api-keys 確認 Key 狀態"
                    )
                elif '429' in str(e):
                    raise ValueError(
                        "OpenAI API 請求頻率限制。\n\n"
                        "請稍後再試或升級帳戶。"
                    )
                elif 'billing' in error_msg.lower() or 'quota' in error_msg.lower():
                    raise ValueError(
                        "OpenAI 帳戶問題。\n\n"
                        "可能原因：\n"
                        "1. 未綁定付款方式\n"
                        "2. 餘額不足\n"
                        "3. 配額用完\n\n"
                        "請訪問：https://platform.openai.com/account/billing"
                    )
                else:
                    raise ValueError(f"OpenAI API 錯誤：{error_msg}")
            except ValueError:
                raise
            except:
                raise ValueError(f"OpenAI API 錯誤：{str(e)}")
        
        except Exception as e:
            raise ValueError(f"API 調用失敗：{str(e)}")
    
    def generate_document(self, doc_type: str, template_path: str, 
                         user_requirements: str, output_format: str = 'docx') -> dict:
        """生成文檔"""
        try:
            logger.info(f"開始生成文檔 - 類型: {doc_type}, 格式: {output_format}")
            
            # 提取模板內容
            template_content = FileProcessor.extract_text(template_path)
            if not template_content or len(template_content.strip()) < 10:
                return {
                    "success": False,
                    "error": "無法提取模板內容或內容過短"
                }
            
            # 獲取文檔配置
            doc_config = DOCUMENT_TYPES.get(doc_type)
            if not doc_config:
                return {
                    "success": False,
                    "error": f"未知的文檔類型: {doc_type}"
                }
            
            # 構建提示詞（限制模板長度）
            prompt = doc_config['prompt_template'].format(
                template_content=template_content[:5000],  # 限制模板內容長度
                user_requirements=user_requirements if user_requirements else "請根據模板生成專業的文檔內容"
            )
            
            # 調用 AI API 生成內容
            logger.info(f"使用 {self.api_type} API 生成內容...")
            if self.api_type == 'openai':
                generated_content = self.call_openai_api(prompt)
            else:
                generated_content = self.call_gemini_api(prompt)
            
            # 生成輸出文件
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_filename = f"{doc_type}_{timestamp}.{output_format}"
            output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
            
            logger.info(f"轉換為 {output_format} 格式...")
            
            if output_format == 'docx':
                output_doc = FormatConverter.markdown_to_docx(generated_content, doc_config)
                output_doc.save(output_path)
            
            elif output_format == 'pptx':
                output_pptx = FormatConverter.markdown_to_pptx(generated_content, doc_config)
                output_pptx.save(output_path)
            
            elif output_format == 'pdf':
                FormatConverter.markdown_to_pdf(generated_content, doc_config, output_path)
            
            elif output_format == 'md':
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(f"# {doc_config['name']}\n\n")
                    f.write(f"**生成時間：** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                    f.write(f"---\n\n")
                    f.write(generated_content)
            
            logger.info(f"✅ 文檔生成成功: {output_filename}")
            
            return {
                "success": True,
                "filename": output_filename,
                "path": output_path,
                "format": output_format,
                "content_preview": generated_content[:1000] + "..." if len(generated_content) > 1000 else generated_content,
                "api_type": self.api_type,
                "model": self.api_config.get('openai_model') if self.api_type == 'openai' else 'gemini-2.0-flash-exp'
            }
            
        except Exception as e:
            logger.error(f"❌ 生成文檔失敗: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }


# 初始化文檔生成器
doc_generator = DocumentGenerator()


def allowed_file(filename):
    """檢查文件類型是否允許"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


# ==================== API 路由 ====================

@app.route('/')
def index():
    """首頁"""
    return render_template('index_v3.html')


@app.route('/api/config', methods=['GET', 'POST'])
def api_config():
    """API 配置"""
    if request.method == 'GET':
        config = doc_generator.api_config
        return jsonify({
            "api_type": config.get('api_type', 'gemini'),
            "has_gemini_key": bool(config.get('gemini_api_key')),
            "has_openai_key": bool(config.get('openai_api_key')),
            "openai_model": config.get('openai_model', 'gpt-4o-mini')
        })
    
    elif request.method == 'POST':
        try:
            data = request.json
            doc_generator.save_api_config(data)
            return jsonify({"success": True, "message": "配置已保存"})
        except Exception as e:
            return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/upload_template', methods=['POST'])
def upload_template():
    """上傳模板"""
    if 'file' not in request.files:
        return jsonify({"success": False, "error": "沒有文件"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"success": False, "error": "沒有選擇文件"}), 400
    
    if file and allowed_file(file.filename):
        try:
            filename = secure_filename(file.filename)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            name, ext = os.path.splitext(filename)
            filename = f"{timestamp}_{name}{ext}"
            filepath = os.path.join(app.config['TEMPLATES_FOLDER'], filename)
            file.save(filepath)
            
            # 提取預覽
            preview = FileProcessor.extract_text(filepath)[:200] + "..."
            
            logger.info(f"模板上傳成功: {filename}")
            
            return jsonify({
                "success": True,
                "filename": filename,
                "path": filepath,
                "preview": preview,
                "message": "模板上傳成功"
            })
        except Exception as e:
            logger.error(f"模板上傳失敗: {e}")
            return jsonify({"success": False, "error": str(e)}), 500
    
    return jsonify({"success": False, "error": "不支持的文件格式"}), 400


@app.route('/api/templates', methods=['GET'])
def list_templates():
    """列出所有模板"""
    templates = []
    template_dir = Path(app.config['TEMPLATES_FOLDER'])
    
    if template_dir.exists():
        for file_path in template_dir.glob('*'):
            if file_path.suffix.lower()[1:] in ALLOWED_EXTENSIONS:
                templates.append({
                    "filename": file_path.name,
                    "type": file_path.suffix.upper()[1:],
                    "size": file_path.stat().st_size,
                    "modified": datetime.fromtimestamp(file_path.stat().st_mtime).strftime('%Y-%m-%d %H:%M:%S')
                })
    
    return jsonify(templates)


@app.route('/api/view_template/<filename>', methods=['GET'])
def view_template(filename):
    """檢視模板內容"""
    try:
        file_path = os.path.join(app.config['TEMPLATES_FOLDER'], filename)
        if not os.path.exists(file_path):
            return jsonify({"success": False, "error": "模板文件不存在"}), 404
        
        # 提取模板內容
        content = FileProcessor.extract_text(file_path)
        
        return jsonify({
            "success": True,
            "filename": filename,
            "content": content
        })
    except Exception as e:
        logger.error(f"檢視模板失敗: {e}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/document_types', methods=['GET'])
def get_document_types():
    """獲取文檔類型列表"""
    types = []
    for key, value in DOCUMENT_TYPES.items():
        types.append({
            "id": key,
            "name": value['name'],
            "description": value['description']
        })
    return jsonify(types)


@app.route('/api/output_formats', methods=['GET'])
def get_output_formats():
    """獲取輸出格式列表"""
    return jsonify([
        {"id": "docx", "name": "Word文檔 (.docx)", "icon": "📄"},
        {"id": "pptx", "name": "PowerPoint簡報 (.pptx)", "icon": "📊"},
        {"id": "pdf", "name": "PDF文檔 (.pdf)", "icon": "📕"},
        {"id": "md", "name": "Markdown (.md)", "icon": "📝"}
    ])


@app.route('/api/generate', methods=['POST'])
def generate_document():
    """生成文檔"""
    try:
        data = request.json
        doc_type = data.get('doc_type')
        template_filename = data.get('template')
        user_requirements = data.get('requirements', '')
        output_format = data.get('output_format', 'docx')
        
        # 驗證參數
        if not doc_type:
            return jsonify({"success": False, "error": "請選擇文檔類型"}), 400
        
        if not template_filename:
            return jsonify({"success": False, "error": "請選擇模板"}), 400
        
        if output_format not in OUTPUT_FORMATS:
            return jsonify({"success": False, "error": "不支持的輸出格式"}), 400
        
        # 檢查模板文件
        template_path = os.path.join(app.config['TEMPLATES_FOLDER'], template_filename)
        if not os.path.exists(template_path):
            return jsonify({"success": False, "error": "模板文件不存在"}), 404
        
        # 生成文檔
        result = doc_generator.generate_document(doc_type, template_path, user_requirements, output_format)
        
        if result['success']:
            return jsonify(result)
        else:
            return jsonify(result), 500
    
    except Exception as e:
        logger.error(f"生成文檔時發生錯誤: {e}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/download/<filename>')
def download_document(filename):
    """下載文檔"""
    file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True, download_name=filename)
    else:
        return jsonify({"error": "文件不存在"}), 404


@app.route('/api/generated_documents', methods=['GET'])
def list_generated_documents():
    """列出已生成的文檔"""
    documents = []
    output_dir = Path(app.config['OUTPUT_FOLDER'])
    
    if output_dir.exists():
        for file_path in output_dir.glob('*'):
            if file_path.suffix[1:] in OUTPUT_FORMATS:
                documents.append({
                    "filename": file_path.name,
                    "format": file_path.suffix[1:].upper(),
                    "size": file_path.stat().st_size,
                    "created": datetime.fromtimestamp(file_path.stat().st_ctime).strftime('%Y-%m-%d %H:%M:%S')
                })
    
    documents.sort(key=lambda x: x['created'], reverse=True)
    return jsonify(documents)


@app.route('/api/delete_template/<filename>', methods=['DELETE'])
def delete_template(filename):
    """刪除模板"""
    try:
        file_path = os.path.join(app.config['TEMPLATES_FOLDER'], filename)
        if os.path.exists(file_path):
            os.remove(file_path)
            logger.info(f"模板已刪除: {filename}")
            return jsonify({"success": True, "message": "模板已刪除"})
        else:
            return jsonify({"success": False, "error": "文件不存在"}), 404
    except Exception as e:
        logger.error(f"刪除模板失敗: {e}")
        return jsonify({"success": False, "error": str(e)}), 500


# ==================== 應用啟動 ====================

if __name__ == '__main__':
    logger.info("=" * 60)
    logger.info("AI文檔生成器 V3.0.3 啟動中...")
    logger.info("=" * 60)
    
    print("\n" + "=" * 60)
    print("🤖 AI文檔生成器 V3.0.3 - 終極版（含成本追蹤）")
    print("=" * 60)
    print(f"支持輸入格式: DOCX, PPTX, TXT, MD" + (" PDF," if PYMUPDF_AVAILABLE else " (PDF讀取不可用)"))
    print(f"支持輸出格式: DOCX, PPTX, PDF, Markdown")
    print(f"AI 引擎: Google Gemini 2.0 / OpenAI GPT (帶成本追蹤)")
    print(f"Web 界面: http://localhost:5000")
    print("=" * 60)
    
    if not PYMUPDF_AVAILABLE:
        print("\n⚠️  提示：PDF 模板讀取功能不可用")
        print("   建議使用 Word (.docx) 或文本 (.txt) 格式的模板")
    
    print("\n💰 成本追蹤:")
    print("   - OpenAI API 調用會顯示實時成本")
    print("   - 成本記錄保存在: config/cost_log.csv")
    print("   - 查看統計: python view_costs.py")
    
    print("\n" + "=" * 60 + "\n")
    
    app.run(debug=True, host='0.0.0.0', port=5000)